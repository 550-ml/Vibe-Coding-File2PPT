from __future__ import annotations

import math
from collections.abc import Callable
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt

from .models import BuildResult, ProjectTree, Topic

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def build_ppt(
    project_tree: ProjectTree,
    output_path: str | Path,
    progress_callback: Callable[[int, int, str], None] | None = None,
) -> BuildResult:
    output = Path(output_path).expanduser().resolve()
    output.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT

    tree_pages = _build_directory_tree_pages(project_tree)
    total_slides = _estimate_total_slides(project_tree, len(tree_pages))
    built_slides = 0

    _add_cover_slide(presentation, project_tree.root_name)
    built_slides += 1
    if progress_callback is not None:
        progress_callback(built_slides, total_slides, project_tree.root_name)
    built_slides += _add_directory_tree_slides(presentation, tree_pages)
    if progress_callback is not None:
        progress_callback(built_slides, total_slides, "目录结构")

    for section in project_tree.sections:
        for topic in section.topics:
            added = _add_topic_slides(presentation, section.name, topic)
            built_slides += added
            if progress_callback is not None:
                progress_callback(built_slides, total_slides, f"{section.name} -- {topic.name}")

    presentation.save(output)
    return BuildResult(
        output_path=output,
        slide_count=len(presentation.slides),
        skipped_files=project_tree.skipped_files,
    )


def _estimate_total_slides(project_tree: ProjectTree, tree_slide_count: int) -> int:
    total = 1 + tree_slide_count
    for section in project_tree.sections:
        for topic in section.topics:
            total += math.ceil(len(topic.items) / _items_per_slide(len(topic.items)))
    return total


def _add_cover_slide(presentation: Presentation, root_name: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _apply_slide_background(slide, RGBColor(244, 247, 252))

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.0),
        Inches(0.0),
        Inches(13.333),
        Inches(0.42),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(33, 84, 166)
    accent.line.fill.background()

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.95),
        Inches(1.2),
        Inches(11.4),
        Inches(4.8),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(218, 226, 236)

    box = slide.shapes.add_textbox(Inches(1.35), Inches(2.0), Inches(10.4), Inches(1.2))
    frame = box.text_frame
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = root_name
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = RGBColor(28, 42, 60)

    sub = slide.shapes.add_textbox(Inches(1.38), Inches(3.15), Inches(6.8), Inches(1.0))
    sub_frame = sub.text_frame
    sub_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    sub_p = sub_frame.paragraphs[0]
    sub_p.alignment = PP_ALIGN.LEFT
    sub_run = sub_p.add_run()
    sub_run.text = "目录结构总览与文件缩略图自动生成"
    sub_run.font.name = "Microsoft YaHei"
    sub_run.font.size = Pt(15)
    sub_run.font.color.rgb = RGBColor(84, 98, 115)

    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(1.38),
        Inches(4.35),
        Inches(2.35),
        Inches(0.48),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(229, 238, 251)
    badge.line.fill.background()
    badge_frame = badge.text_frame
    badge_p = badge_frame.paragraphs[0]
    badge_p.alignment = PP_ALIGN.CENTER
    badge_run = badge_p.add_run()
    badge_run.text = "Folder2PPT"
    badge_run.font.name = "Microsoft YaHei"
    badge_run.font.size = Pt(12)
    badge_run.font.bold = True
    badge_run.font.color.rgb = RGBColor(33, 84, 166)


def _add_directory_tree_slides(
    presentation: Presentation,
    tree_pages: list[tuple[str, str, list[str]]],
) -> int:
    for index, (title, subtitle, lines) in enumerate(tree_pages, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _apply_slide_background(slide, RGBColor(248, 250, 253))
        page_title = title
        if len(tree_pages) > 1:
            page_title = f"{title} ({index}/{len(tree_pages)})"
        _add_page_title(slide, page_title, subtitle=subtitle, badge="目录")

        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.72),
            Inches(1.45),
            Inches(11.85),
            Inches(5.5),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(220, 228, 238)

        box = slide.shapes.add_textbox(Inches(0.95), Inches(1.7), Inches(11.35), Inches(4.95))
        frame = box.text_frame
        frame.word_wrap = False
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = "\n".join(lines)
        run.font.name = "Consolas"
        run.font.size = Pt(_tree_font_size(len(lines)))
        run.font.color.rgb = RGBColor(46, 58, 72)
    return len(tree_pages)


def _build_directory_tree_pages(project_tree: ProjectTree) -> list[tuple[str, str, list[str]]]:
    lines = _build_directory_tree_lines(project_tree)
    stats = _build_tree_stats(project_tree)
    total_lines = len(lines)

    if total_lines <= 26:
        return [("目录总览", stats, lines)]

    if total_lines <= 34:
        return [("目录总览", f"{stats} | 已压缩展示", lines)]

    chunks = _paginate_tree_lines(lines, 30)
    pages: list[tuple[str, str, list[str]]] = []
    for index, chunk in enumerate(chunks, start=1):
        subtitle = stats if index == 1 else f"{stats} | 续页 {index}/{len(chunks)}"
        pages.append((f"目录总览 {index}/{len(chunks)}", subtitle, chunk))
    return pages


def _build_directory_tree_lines(project_tree: ProjectTree) -> list[str]:
    tree: dict[str, dict] = {}
    for section in project_tree.sections:
        for topic in section.topics:
            for item in topic.items:
                parts = item.source_path.relative_to(project_tree.root_path).parts
                node = tree
                for part in parts[:-1]:
                    node = node.setdefault(part, {})
                node[item.display_name] = None

    lines = [f"{project_tree.root_name}/"]
    lines.extend(_render_tree_nodes(tree, prefix=""))
    return lines


def _build_tree_stats(project_tree: ProjectTree) -> str:
    section_count = len(project_tree.sections)
    file_count = sum(len(topic.items) for section in project_tree.sections for topic in section.topics)
    years: set[str] = set()
    for section in project_tree.sections:
        for topic in section.topics:
            for item in topic.items:
                for part in item.source_path.relative_to(project_tree.root_path).parts:
                    if part.isdigit() and len(part) == 4:
                        years.add(part)
    year_text = "未识别年份"
    if years:
        ordered_years = sorted(years)
        year_text = ordered_years[0] if len(ordered_years) == 1 else f"{ordered_years[0]}-{ordered_years[-1]}"
    return f"统计：{section_count} 个领域 | {file_count} 篇文件 | {year_text}"


def _render_tree_nodes(tree: dict[str, dict | None], prefix: str) -> list[str]:
    lines: list[str] = []
    items = list(tree.items())
    for index, (name, child) in enumerate(items):
        is_last = index == len(items) - 1
        branch = "└── " if is_last else "├── "
        suffix = "/" if isinstance(child, dict) else ""
        lines.append(f"{prefix}{branch}{name}{suffix}")
        if isinstance(child, dict):
            child_prefix = f"{prefix}{'    ' if is_last else '│   '}"
            lines.extend(_render_tree_nodes(child, child_prefix))
    return lines


def _paginate_tree_lines(lines: list[str], max_lines_per_slide: int = 24) -> list[list[str]]:
    return [
        lines[index : index + max_lines_per_slide]
        for index in range(0, len(lines), max_lines_per_slide)
    ]


def _tree_font_size(line_count: int) -> float:
    if line_count <= 22:
        return 11.8
    if line_count <= 28:
        return 10.8
    return 9.8


def _add_topic_slides(presentation: Presentation, section_name: str, topic: Topic) -> int:
    items = topic.items
    per_slide = _items_per_slide(len(items))
    total_pages = math.ceil(len(items) / per_slide)
    for page_index in range(total_pages):
        chunk = items[page_index * per_slide : (page_index + 1) * per_slide]
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _apply_slide_background(slide, RGBColor(248, 250, 253))
        title = f"{section_name} -- {topic.name}"
        if total_pages > 1:
            title = f"{title} ({page_index + 1}/{total_pages})"
        _add_page_title(
            slide,
            title,
            subtitle=f"本页 {len(chunk)} 项，共 {len(items)} 项",
            badge=section_name,
        )
        _add_items_grid(slide, chunk)
    return total_pages


def _items_per_slide(item_count: int) -> int:
    if item_count <= 4:
        return 4
    if item_count <= 6:
        return 6
    return 8


def _apply_slide_background(slide, color: RGBColor) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = color


def _add_page_title(slide, title: str, subtitle: str | None = None, badge: str | None = None) -> None:
    if badge:
        badge_shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.72),
            Inches(0.32),
            Inches(1.55),
            Inches(0.36),
        )
        badge_shape.fill.solid()
        badge_shape.fill.fore_color.rgb = RGBColor(229, 238, 251)
        badge_shape.line.fill.background()
        badge_frame = badge_shape.text_frame
        badge_p = badge_frame.paragraphs[0]
        badge_p.alignment = PP_ALIGN.CENTER
        badge_run = badge_p.add_run()
        badge_run.text = badge
        badge_run.font.name = "Microsoft YaHei"
        badge_run.font.size = Pt(11)
        badge_run.font.bold = True
        badge_run.font.color.rgb = RGBColor(33, 84, 166)

    shape = slide.shapes.add_textbox(Inches(0.72), Inches(0.72), Inches(12), Inches(0.52))
    frame = shape.text_frame
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = RGBColor(31, 44, 58)

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.74), Inches(1.13), Inches(10.8), Inches(0.35))
        subtitle_frame = subtitle_box.text_frame
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_run = subtitle_p.add_run()
        subtitle_run.text = subtitle
        subtitle_run.font.name = "Microsoft YaHei"
        subtitle_run.font.size = Pt(11)
        subtitle_run.font.color.rgb = RGBColor(102, 115, 129)

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.72),
        Inches(1.48),
        Inches(11.95),
        Inches(0.03),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(210, 220, 233)
    line.line.fill.background()


def _add_items_grid(slide, items) -> None:
    count = len(items)
    columns = 2 if count <= 4 else (3 if count <= 6 else 4)
    box_width = 11.85 / columns
    card_width = box_width - 0.18
    image_width = card_width - 0.18
    image_height = 2.0 if columns == 4 else 2.25
    rows = math.ceil(count / columns)
    start_left = 0.7
    start_top = 1.78
    row_step = 2.92 if rows <= 2 else 2.66

    for index, item in enumerate(items):
        row = index // columns
        col = index % columns
        left = Inches(start_left + col * box_width)
        top = Inches(start_top + row * row_step)

        if item.preview_path is None:
            raise ValueError(f"缺少预览图: {item.source_path}")

        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            top,
            Inches(card_width),
            Inches(image_height + 0.6),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(220, 228, 238)

        slide.shapes.add_picture(
            str(item.preview_path),
            left + Inches(0.09),
            top + Inches(0.08),
            width=Inches(image_width),
            height=Inches(image_height),
        )

        caption = slide.shapes.add_textbox(
            left + Inches(0.06),
            top + Inches(image_height + 0.14),
            Inches(card_width - 0.12),
            Inches(0.46),
        )
        frame = caption.text_frame
        frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = item.display_name
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(11 if columns == 4 else 12)
        run.font.color.rgb = RGBColor(54, 64, 74)
